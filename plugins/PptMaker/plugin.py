"""
PptMaker: built-in plugin to create PowerPoint (.pptx) presentations.
Runs with Core; requires python-pptx.
"""
import json
import os
import re
from pathlib import Path
from typing import Any, Dict, List

from loguru import logger

from base.BasePlugin import BasePlugin
from base.util import Util
from base.workspace import get_workspace_dir
from core.coreInterface import CoreInterface

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

LAYOUT_TITLE = 0
LAYOUT_TITLE_BODY = 1


def _add_bullets(shape, bullets: List[str], level: int = 0) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = level
        p.space_after = Pt(6)


def _build_pptx(main_title: str, subtitle: str, slides_list: List[Dict[str, Any]], out_path: Path) -> None:
    """Build a presentation and save to out_path. slides_list: [{"title": str, "bullets": [str]}, ...]."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Install it with: pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    slide.shapes.title.text = main_title or "Presentation"
    if subtitle and hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    # Content slides
    for item in slides_list:
        title = (item.get("title") or "").strip() or "Slide"
        bullets = item.get("bullets")
        if not isinstance(bullets, list):
            bullets = [str(bullets)] if bullets else []
        bullets = [str(b).strip() for b in bullets if b]

        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
        slide.shapes.title.text = title
        if bullets:
            _add_bullets(slide.placeholders[1], bullets)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _parse_outline(outline: str) -> tuple:
    """Parse markdown-style outline into (main_title, subtitle, slides_list).
    ## starts a new slide (title = rest of line). - or * are bullets.
    First line or first block = title slide (first line = main_title, optional second = subtitle).
    """
    main_title = ""
    subtitle = ""
    slides_list: List[Dict[str, Any]] = []
    current_title = None
    current_bullets: List[str] = []

    def flush_slide():
        nonlocal current_title, current_bullets
        if current_title is not None:
            slides_list.append({"title": current_title, "bullets": list(current_bullets)})
        current_title = None
        current_bullets = []

    lines = (outline or "").strip().split("\n")
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("##"):
            flush_slide()
            current_title = stripped[2:].strip() or "Slide"
            current_bullets = []
        elif stripped.startswith("-") or stripped.startswith("*"):
            bullet = stripped[1:].strip()
            if current_title is not None:
                current_bullets.append(bullet)
            else:
                if not main_title:
                    main_title = bullet
                elif not subtitle:
                    subtitle = bullet
                else:
                    if not current_title:
                        current_title = "Content"
                    current_bullets.append(bullet)
        elif stripped:
            if current_title is not None:
                current_bullets.append(stripped)
            else:
                if not main_title:
                    main_title = stripped
                elif not subtitle:
                    subtitle = stripped
                else:
                    current_title = "Content"
                    current_bullets.append(stripped)
    flush_slide()

    if not main_title and slides_list:
        main_title = slides_list[0].get("title") or "Presentation"
    return main_title, subtitle, slides_list


def _parse_document_to_slides(content: str, doc_title: str = "") -> List[Dict[str, Any]]:
    """Parse document text into slides. Markdown: ## = slide title, - or * = bullets. Plain text: split by \\n\\n or by ##; first line = title, rest = bullets. Long bullets are truncated per slide."""
    if not (content or "").strip():
        return [{"title": doc_title or "Document", "bullets": ["(No content)"]}]
    text = (content or "").strip()
    slides: List[Dict[str, Any]] = []
    # If it looks like markdown (has ##), use same logic as outline
    if "##" in text or text.startswith("#"):
        main_title, subtitle, slides_list = _parse_outline(text)
        if slides_list:
            for s in slides_list:
                slides.append(s)
        elif main_title:
            slides.append({"title": main_title, "bullets": [subtitle] if subtitle else []})
    else:
        # Plain text: split into blocks (double newline or long lines), each block = slide or bullet set
        blocks = re.split(r"\n\s*\n", text)
        max_bullets_per_slide = 12
        max_bullet_len = 400
        for i, block in enumerate(blocks):
            block = block.strip()
            if not block:
                continue
            lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
            if not lines:
                continue
            title = lines[0][:120]
            bullets = []
            for ln in lines[1:]:
                bullets.append(ln[:max_bullet_len] if len(ln) > max_bullet_len else ln)
            if not bullets and len(lines) == 1 and len(lines[0]) > 80:
                # Single long line: treat as one bullet or split by sentence
                first = lines[0]
                for part in re.split(r"(?<=[.!?])\s+", first)[:max_bullets_per_slide]:
                    if part.strip():
                        bullets.append(part.strip()[:max_bullet_len])
            if len(bullets) > max_bullets_per_slide:
                slides.append({"title": title, "bullets": bullets[:max_bullets_per_slide]})
                # Overflow into next slide
                for j in range(max_bullets_per_slide, len(bullets), max_bullets_per_slide):
                    chunk = bullets[j : j + max_bullets_per_slide]
                    slides.append({"title": f"{title} (continued)", "bullets": chunk})
            else:
                slides.append({"title": title, "bullets": bullets or [title]})
    if not slides and doc_title:
        slides.append({"title": doc_title, "bullets": [text[:500]]})
    return slides


def _resolve_document_path(path_str: str, base_dirs: List[Path]) -> Path:
    """Resolve a path against allowed base dirs (workspace, project root). Return Path if safe and file exists, else None."""
    path_str = (path_str or "").strip()
    if not path_str:
        return None
    p = Path(path_str)
    bases = [b.resolve() for b in base_dirs]
    if p.is_absolute():
        resolved = p.resolve()
        for base in bases:
            try:
                resolved.relative_to(base)
                if resolved.is_file():
                    return resolved
            except ValueError:
                continue
        return None
    for base in bases:
        candidate = (base / p).resolve()
        try:
            candidate.relative_to(base)
        except ValueError:
            continue
        if candidate.is_file():
            return candidate
    return None


class PptMakerPlugin(BasePlugin):
    def __init__(self, coreInst: CoreInterface):
        super().__init__(coreInst=coreInst)
        self.config = {}
        config_path = Path(__file__).resolve().parent / "config.yml"
        if config_path.exists():
            self.config = Util().load_yml_config(str(config_path)) or {}
        logger.debug("PptMaker plugin config: {}", self.config)

    def _output_dir(self) -> Path:
        out = (self.config.get("output_dir") or "").strip()
        if out:
            p = Path(out)
            if not p.is_absolute():
                p = Path(__file__).resolve().parent.parent.parent / p
            return p
        meta = Util().get_core_metadata()
        ws = get_workspace_dir(getattr(meta, "workspace_dir", None) or "config/workspace")
        return ws / "presentations"

    def _safe_filename(self, name: str) -> str:
        name = re.sub(r'[^\w\s\-\.]', '', name)
        name = name.strip()[:80] or "presentation"
        return name + ".pptx" if not name.lower().endswith(".pptx") else name

    async def create_presentation(self) -> str:
        """Create .pptx from structured params (main_title, subtitle, slides JSON)."""
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        main_title = (self.config.get("main_title") or "").strip() or "Presentation"
        subtitle = (self.config.get("subtitle") or "").strip()
        slides_raw = self.config.get("slides")
        if isinstance(slides_raw, str):
            try:
                slides_list = json.loads(slides_raw)
            except json.JSONDecodeError:
                return json.dumps({"success": False, "error": "Invalid JSON in slides parameter."})
        elif isinstance(slides_raw, list):
            slides_list = slides_raw
        else:
            return json.dumps({"success": False, "error": "slides must be a JSON array of {title, bullets} objects."})

        if not isinstance(slides_list, list):
            return json.dumps({"success": False, "error": "slides must be an array."})

        out_dir = self._output_dir()
        out_filename = (self.config.get("output_filename") or "").strip()
        if not out_filename:
            from datetime import datetime
            out_filename = self._safe_filename(main_title) or f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        else:
            if not out_filename.lower().endswith(".pptx"):
                out_filename += ".pptx"
        out_path = out_dir / out_filename

        try:
            _build_pptx(main_title, subtitle, slides_list, out_path)
            return json.dumps({
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation saved to {out_path.resolve()}",
            })
        except Exception as e:
            logger.exception(e)
            return json.dumps({"success": False, "error": str(e)})

    async def create_from_outline(self) -> str:
        """Create .pptx from markdown-style outline (## titles, - bullets)."""
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        outline = (self.config.get("outline") or "").strip()
        if not outline:
            return json.dumps({"success": False, "error": "outline parameter is required."})

        main_title, subtitle, slides_list = _parse_outline(outline)
        if not slides_list and not main_title:
            return json.dumps({"success": False, "error": "Could not parse any slides from the outline."})

        out_dir = self._output_dir()
        out_filename = (self.config.get("output_filename") or "").strip()
        if not out_filename:
            from datetime import datetime
            out_filename = self._safe_filename(main_title) or f"outline_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        else:
            if not out_filename.lower().endswith(".pptx"):
                out_filename += ".pptx"
        out_path = out_dir / out_filename

        try:
            _build_pptx(main_title, subtitle, slides_list, out_path)
            return json.dumps({
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation created from outline and saved to {out_path.resolve()}",
            })
        except Exception as e:
            logger.exception(e)
            return json.dumps({"success": False, "error": str(e)})

    async def create_from_documents(self) -> str:
        """Create .pptx from one or more documents (file paths or pre-read content)."""
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        paths_raw = self.config.get("document_paths")
        contents_raw = self.config.get("document_contents")
        if not paths_raw and not contents_raw:
            return json.dumps({"success": False, "error": "Provide at least one of document_paths or document_contents."})

        main_title = (self.config.get("main_title") or "").strip()
        all_slides: List[Dict[str, Any]] = []
        sources: List[str] = []

        if paths_raw:
            if isinstance(paths_raw, str):
                try:
                    path_list = json.loads(paths_raw)
                except json.JSONDecodeError:
                    return json.dumps({"success": False, "error": "document_paths must be a JSON array of strings."})
            else:
                path_list = paths_raw if isinstance(paths_raw, list) else []
            if not isinstance(path_list, list):
                return json.dumps({"success": False, "error": "document_paths must be an array."})
            meta = Util().get_core_metadata()
            ws = get_workspace_dir(getattr(meta, "workspace_dir", None) or "config/workspace")
            project_root = Path(__file__).resolve().parent.parent.parent
            base_dirs = [ws, project_root]
            for path_str in path_list:
                if not isinstance(path_str, str):
                    continue
                resolved = _resolve_document_path(path_str, base_dirs)
                if not resolved:
                    return json.dumps({"success": False, "error": f"File not found or not allowed: {path_str}"})
                try:
                    content = resolved.read_text(encoding="utf-8", errors="replace")
                except Exception as e:
                    return json.dumps({"success": False, "error": f"Could not read {path_str}: {e}"})
                doc_title = resolved.stem
                slides = _parse_document_to_slides(content, doc_title)
                if not main_title and slides:
                    main_title = main_title or doc_title or slides[0].get("title", "Document")
                all_slides.extend(slides)
                sources.append(resolved.name)

        if contents_raw:
            if isinstance(contents_raw, str):
                try:
                    content_list = json.loads(contents_raw)
                except json.JSONDecodeError:
                    return json.dumps({"success": False, "error": "document_contents must be a JSON array of {title, content} objects."})
            else:
                content_list = contents_raw if isinstance(contents_raw, list) else []
            if not isinstance(content_list, list):
                return json.dumps({"success": False, "error": "document_contents must be an array."})
            for i, item in enumerate(content_list):
                if not isinstance(item, dict):
                    continue
                title = (item.get("title") or item.get("name") or f"Document {i+1}").strip()
                content = (item.get("content") or "").strip()
                slides = _parse_document_to_slides(content, title)
                if not main_title and slides:
                    main_title = main_title or title
                all_slides.extend(slides)
                sources.append(title)

        if not all_slides:
            return json.dumps({"success": False, "error": "No slides could be parsed from the document(s)."})
        if not main_title:
            main_title = all_slides[0].get("title", "Presentation") if all_slides else "Presentation"

        out_dir = self._output_dir()
        out_filename = (self.config.get("output_filename") or "").strip()
        if not out_filename:
            from datetime import datetime
            out_filename = self._safe_filename(main_title) or f"documents_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        else:
            if not out_filename.lower().endswith(".pptx"):
                out_filename += ".pptx"
        out_path = out_dir / out_filename

        try:
            _build_pptx(main_title, "", all_slides, out_path)
            return json.dumps({
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation created from {len(sources)} document(s) and saved to {out_path.resolve()}",
                "sources": sources,
            })
        except Exception as e:
            logger.exception(e)
            return json.dumps({"success": False, "error": str(e)})

    async def run(self):
        """Default: create_from_outline if user_input looks like outline, else suggest capabilities."""
        text = (getattr(self.promptRequest, "text", None) or self.user_input or "").strip()
        if "##" in text or (text.startswith("-") or "\n-" in text):
            self.config["outline"] = text
            return await self.create_from_outline()
        return "Use capability create_presentation (with main_title and slides JSON) or create_from_outline (with outline text) to generate a PPT."

    def initialize(self):
        if self.initialized:
            return
        logger.debug("Initializing PptMaker plugin")
        super().initialize()
        self.initialized = True
