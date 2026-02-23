"""
PPT Generation plugin: create PowerPoint (.pptx) presentations.
Parameters (title, subtitle, slides) can be extracted from the source (web search results,
other files, or given text) — use create_from_source to pass raw source, or
create_from_documents/create_from_outline with content from web_search/file_read.
Output file is saved to the user's private output folder (base/{user_id}/output/ or
companion/output/) and the response includes an open link. Requires python-pptx.
"""
import json
import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List

from base.BasePlugin import BasePlugin
from base.util import Util
from base.workspace import get_workspace_dir
from core.coreInterface import CoreInterface

try:
    from loguru import logger
except ImportError:
    logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

LAYOUT_TITLE = 0
LAYOUT_TITLE_BODY = 1


def _add_bullets(shape, bullets: List[str], level: int = 0) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = level
        p.space_after = Pt(6)


def _build_pptx(main_title: str, subtitle: str, slides_list: List[Dict[str, Any]], out_path: Path) -> None:
    """Build a presentation and save to out_path. slides_list: [{"title": str, "bullets": [str]}, ...]."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Install it with: pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    slide.shapes.title.text = main_title or "Presentation"
    if subtitle and hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    for item in slides_list:
        title = (item.get("title") or "").strip() or "Slide"
        bullets = item.get("bullets")
        if not isinstance(bullets, list):
            bullets = [str(bullets)] if bullets else []
        bullets = [str(b).strip() for b in bullets if b]

        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
        slide.shapes.title.text = title
        if bullets:
            _add_bullets(slide.placeholders[1], bullets)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _parse_outline(outline: str) -> tuple:
    main_title = ""
    subtitle = ""
    slides_list: List[Dict[str, Any]] = []
    current_title = None
    current_bullets: List[str] = []

    def flush_slide():
        nonlocal current_title, current_bullets
        if current_title is not None:
            slides_list.append({"title": current_title, "bullets": list(current_bullets)})
        current_title = None
        current_bullets = []

    lines = (outline or "").strip().split("\n")
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("##"):
            flush_slide()
            current_title = stripped[2:].strip() or "Slide"
            current_bullets = []
        elif stripped.startswith("-") or stripped.startswith("*"):
            bullet = stripped[1:].strip()
            if current_title is not None:
                current_bullets.append(bullet)
            else:
                if not main_title:
                    main_title = bullet
                elif not subtitle:
                    subtitle = bullet
                else:
                    if not current_title:
                        current_title = "Content"
                    current_bullets.append(bullet)
        elif stripped:
            if current_title is not None:
                current_bullets.append(stripped)
            else:
                if not main_title:
                    main_title = stripped
                elif not subtitle:
                    subtitle = stripped
                else:
                    current_title = "Content"
                    current_bullets.append(stripped)
    flush_slide()

    if not main_title and slides_list:
        main_title = slides_list[0].get("title") or "Presentation"
    return main_title, subtitle, slides_list


def _parse_document_to_slides(content: str, doc_title: str = "") -> List[Dict[str, Any]]:
    if not (content or "").strip():
        return [{"title": doc_title or "Document", "bullets": ["(No content)"]}]
    text = (content or "").strip()
    slides: List[Dict[str, Any]] = []
    if "##" in text or text.startswith("#"):
        main_title, subtitle, slides_list = _parse_outline(text)
        if slides_list:
            for s in slides_list:
                slides.append(s)
        elif main_title:
            slides.append({"title": main_title, "bullets": [subtitle] if subtitle else []})
    else:
        blocks = re.split(r"\n\s*\n", text)
        max_bullets_per_slide = 12
        max_bullet_len = 400
        for i, block in enumerate(blocks):
            block = block.strip()
            if not block:
                continue
            lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
            if not lines:
                continue
            title = lines[0][:120]
            bullets = []
            for ln in lines[1:]:
                bullets.append(ln[:max_bullet_len] if len(ln) > max_bullet_len else ln)
            if not bullets and len(lines) == 1 and len(lines[0]) > 80:
                first = lines[0]
                for part in re.split(r"(?<=[.!?])\s+", first)[:max_bullets_per_slide]:
                    if part.strip():
                        bullets.append(part.strip()[:max_bullet_len])
            if len(bullets) > max_bullets_per_slide:
                slides.append({"title": title, "bullets": bullets[:max_bullets_per_slide]})
                for j in range(max_bullets_per_slide, len(bullets), max_bullets_per_slide):
                    chunk = bullets[j : j + max_bullets_per_slide]
                    slides.append({"title": f"{title} (continued)", "bullets": chunk})
            else:
                slides.append({"title": title, "bullets": bullets or [title]})
    if not slides and doc_title:
        slides.append({"title": doc_title, "bullets": [text[:500]]})
    return slides


def _resolve_document_path(path_str: str, base_dirs: List[Path]) -> Path:
    path_str = (path_str or "").strip()
    if not path_str:
        return None
    p = Path(path_str)
    bases = [b.resolve() for b in base_dirs]
    if p.is_absolute():
        resolved = p.resolve()
        for base in bases:
            try:
                resolved.relative_to(base)
                if resolved.is_file():
                    return resolved
            except ValueError:
                continue
        return None
    for base in bases:
        candidate = (base / p).resolve()
        try:
            candidate.relative_to(base)
        except ValueError:
            continue
        if candidate.is_file():
            return candidate
    return None


class PptGenerationPlugin(BasePlugin):
    def __init__(self, coreInst: CoreInterface):
        super().__init__(coreInst=coreInst)
        self.config = {}
        try:
            config_path = Path(__file__).resolve().parent / "config.yml"
            if config_path.exists():
                self.config = Util().load_yml_config(str(config_path)) or {}
        except Exception as e:
            logger.debug("ppt-generation plugin config: {}", e)

    def _output_dir(self) -> Path:
        request_dir = getattr(self, "request_output_dir", None)
        if request_dir is not None and isinstance(request_dir, Path):
            return request_dir
        out = (self.config.get("output_dir") or "").strip()
        if out:
            p = Path(out)
            if not p.is_absolute():
                p = Path(__file__).resolve().parent.parent.parent / p
            return p
        meta = Util().get_core_metadata()
        ws = get_workspace_dir(getattr(meta, "workspace_dir", None) or "config/workspace")
        return ws / "presentations"

    def _safe_filename(self, name: str) -> str:
        name = re.sub(r'[^\w\s\-\.]', '', name)
        name = name.strip()[:80] or "presentation"
        return name + ".pptx" if not name.lower().endswith(".pptx") else name

    async def create_presentation(self) -> str:
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        try:
            main_title = (self.config.get("main_title") or "").strip() or "Presentation"
            subtitle = (self.config.get("subtitle") or "").strip()
            slides_raw = self.config.get("slides")
            if isinstance(slides_raw, str):
                try:
                    slides_list = json.loads(slides_raw)
                except json.JSONDecodeError:
                    return json.dumps({"success": False, "error": "Invalid JSON in slides parameter."})
            elif isinstance(slides_raw, list):
                slides_list = slides_raw
            else:
                return json.dumps({"success": False, "error": "slides must be a JSON array of {title, bullets} objects."})

            if not isinstance(slides_list, list):
                return json.dumps({"success": False, "error": "slides must be an array."})

            out_dir = self._output_dir()
            out_filename = (self.config.get("output_filename") or "").strip()
            if not out_filename:
                from datetime import datetime
                out_filename = self._safe_filename(main_title) or f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            else:
                if not out_filename.lower().endswith(".pptx"):
                    out_filename += ".pptx"
            out_path = out_dir / out_filename
            use_request_output = getattr(self, "request_output_dir", None) is not None
            _build_pptx(main_title, subtitle, slides_list, out_path)
            payload = {
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation saved to {out_path.resolve()}",
            }
            if use_request_output:
                payload["output_rel_path"] = f"output/{out_filename}"
            return json.dumps(payload)
        except Exception as e:
            logger.exception("ppt-generation create_presentation: {}", e)
            return json.dumps({"success": False, "error": str(e)})

    async def create_from_outline(self) -> str:
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        try:
            outline = (self.config.get("outline") or "").strip()
            if not outline:
                return json.dumps({"success": False, "error": "outline parameter is required."})

            main_title, subtitle, slides_list = _parse_outline(outline)
            if not slides_list and not main_title:
                return json.dumps({"success": False, "error": "Could not parse any slides from the outline."})

            out_dir = self._output_dir()
            out_filename = (self.config.get("output_filename") or "").strip()
            if not out_filename:
                from datetime import datetime
                out_filename = self._safe_filename(main_title) or f"outline_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            else:
                if not out_filename.lower().endswith(".pptx"):
                    out_filename += ".pptx"
            out_path = out_dir / out_filename
            use_request_output = getattr(self, "request_output_dir", None) is not None
            _build_pptx(main_title, subtitle, slides_list, out_path)
            payload = {
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation created from outline and saved to {out_path.resolve()}",
            }
            if use_request_output:
                payload["output_rel_path"] = f"output/{out_filename}"
            return json.dumps(payload)
        except Exception as e:
            logger.exception("ppt-generation create_from_outline: {}", e)
            return json.dumps({"success": False, "error": str(e)})

    async def create_from_source(self) -> str:
        """Create a .pptx from a single source: plain text (parsed as outline) or JSON array of {title, content} (e.g. web search results). Extracts main_title and slides from the source. Saves to user's private output folder and returns link."""
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        try:
            source = (self.config.get("source") or "").strip()
            if not source:
                return json.dumps({"success": False, "error": "source parameter is required (plain text outline or JSON array of {title, content})."})
            main_title = ""
            subtitle = ""
            all_slides: List[Dict[str, Any]] = []
            # Try JSON first (e.g. web search results or document_contents)
            stripped = source.strip()
            content_list = None
            if stripped.startswith("["):
                try:
                    content_list = json.loads(source)
                except json.JSONDecodeError:
                    pass
            elif stripped.startswith("{"):
                try:
                    obj = json.loads(source)
                    if isinstance(obj, dict):
                        content_list = obj.get("results") or obj.get("documents") or (obj.get("items") if isinstance(obj.get("items"), list) else None)
                except json.JSONDecodeError:
                    pass
            if isinstance(content_list, list) and content_list:
                for i, item in enumerate(content_list):
                    if not isinstance(item, dict):
                        continue
                    title = str(
                        item.get("title") or item.get("name") or item.get("query") or f"Source {i+1}"
                    ).strip() or f"Source {i+1}"
                    content = str(
                        item.get("content") or item.get("body") or item.get("snippet") or ""
                    ).strip()
                    if not content and isinstance(item.get("results"), list):
                        for r in item.get("results", [])[:5]:
                            if isinstance(r, dict):
                                content += str(r.get("content") or r.get("body") or r.get("snippet") or r.get("title") or "") + "\n"
                    slides = _parse_document_to_slides(content, title)
                    if not main_title and slides:
                        main_title = title or (slides[0].get("title") if slides else "")
                    all_slides.extend(slides)
                if all_slides and not main_title:
                    main_title = all_slides[0].get("title", "Presentation") if all_slides else "Presentation"
            if not all_slides:
                # Plain text: parse as outline
                main_title, subtitle, slides_list = _parse_outline(source)
                if slides_list:
                    all_slides = slides_list
                elif main_title:
                    all_slides = [{"title": main_title, "bullets": [subtitle] if subtitle else []}]
            if not all_slides:
                return json.dumps({"success": False, "error": "Could not extract any slides from the source."})
            if not main_title:
                main_title = (
                    str((all_slides[0] or {}).get("title") or "Presentation")
                    if all_slides
                    else "Presentation"
                )
            main_title = str(main_title or "Presentation").strip() or "Presentation"
            out_dir = self._output_dir()
            out_filename = (self.config.get("output_filename") or "").strip()
            if not out_filename:
                from datetime import datetime
                out_filename = self._safe_filename(main_title) or f"source_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            else:
                if not out_filename.lower().endswith(".pptx"):
                    out_filename += ".pptx"
            out_path = out_dir / out_filename
            use_request_output = getattr(self, "request_output_dir", None) is not None
            _build_pptx(main_title, subtitle or "", all_slides, out_path)
            payload = {
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation created from source and saved to {out_path.resolve()}",
            }
            if use_request_output:
                payload["output_rel_path"] = f"output/{out_filename}"
            return json.dumps(payload)
        except Exception as e:
            logger.exception("ppt-generation create_from_source: {}", e)
            return json.dumps({"success": False, "error": str(e)})

    async def create_from_documents(self) -> str:
        if not _PPTX_AVAILABLE:
            return json.dumps({"success": False, "error": "python-pptx not installed. Install with: pip install python-pptx"})
        try:
            paths_raw = self.config.get("document_paths")
            contents_raw = self.config.get("document_contents")
            if not paths_raw and not contents_raw:
                return json.dumps({"success": False, "error": "Provide at least one of document_paths or document_contents."})

            main_title = (self.config.get("main_title") or "").strip()
            all_slides: List[Dict[str, Any]] = []
            sources: List[str] = []

            if paths_raw:
                if isinstance(paths_raw, str):
                    try:
                        path_list = json.loads(paths_raw)
                    except json.JSONDecodeError:
                        return json.dumps({"success": False, "error": "document_paths must be a JSON array of strings."})
                else:
                    path_list = paths_raw if isinstance(paths_raw, list) else []
                if not isinstance(path_list, list):
                    return json.dumps({"success": False, "error": "document_paths must be an array."})
                meta = Util().get_core_metadata()
                ws = get_workspace_dir(getattr(meta, "workspace_dir", None) or "config/workspace")
                project_root = Path(__file__).resolve().parent.parent.parent
                base_dirs = [ws, project_root]
                for path_str in path_list:
                    if not isinstance(path_str, str):
                        continue
                    resolved = _resolve_document_path(path_str, base_dirs)
                    if not resolved:
                        return json.dumps({"success": False, "error": f"File not found or not allowed: {path_str}"})
                    try:
                        content = resolved.read_text(encoding="utf-8", errors="replace")
                    except Exception as e:
                        return json.dumps({"success": False, "error": f"Could not read {path_str}: {e}"})
                    doc_title = resolved.stem
                    slides = _parse_document_to_slides(content, doc_title)
                    if not main_title and slides:
                        main_title = main_title or doc_title or slides[0].get("title", "Document")
                    all_slides.extend(slides)
                    sources.append(resolved.name)

            if contents_raw:
                if isinstance(contents_raw, str):
                    try:
                        content_list = json.loads(contents_raw)
                    except json.JSONDecodeError:
                        return json.dumps({"success": False, "error": "document_contents must be a JSON array of {title, content} objects."})
                else:
                    content_list = contents_raw if isinstance(contents_raw, list) else []
                if not isinstance(content_list, list):
                    return json.dumps({"success": False, "error": "document_contents must be an array."})
                for i, item in enumerate(content_list):
                    if not isinstance(item, dict):
                        continue
                    title = (item.get("title") or item.get("name") or f"Document {i+1}").strip()
                    content = (item.get("content") or "").strip()
                    slides = _parse_document_to_slides(content, title)
                    if not main_title and slides:
                        main_title = main_title or title
                    all_slides.extend(slides)
                    sources.append(title)

            if not all_slides:
                return json.dumps({"success": False, "error": "No slides could be parsed from the document(s)."})
            if not main_title:
                main_title = all_slides[0].get("title", "Presentation") if all_slides else "Presentation"

            out_dir = self._output_dir()
            out_filename = (self.config.get("output_filename") or "").strip()
            if not out_filename:
                from datetime import datetime
                out_filename = self._safe_filename(main_title) or f"documents_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            else:
                if not out_filename.lower().endswith(".pptx"):
                    out_filename += ".pptx"
            out_path = out_dir / out_filename
            use_request_output = getattr(self, "request_output_dir", None) is not None
            _build_pptx(main_title, "", all_slides, out_path)
            payload = {
                "success": True,
                "path": str(out_path.resolve()),
                "message": f"Presentation created from {len(sources)} document(s) and saved to {out_path.resolve()}",
                "sources": sources,
            }
            if use_request_output:
                payload["output_rel_path"] = f"output/{out_filename}"
            return json.dumps(payload)
        except Exception as e:
            logger.exception("ppt-generation create_from_documents: {}", e)
            return json.dumps({"success": False, "error": str(e)})

    async def run(self) -> str:
        try:
            text = (getattr(self.promptRequest, "text", None) or self.user_input or "").strip()
            if "##" in text or (text.startswith("-") or "\n-" in text):
                self.config["outline"] = text
                return await self.create_from_outline()
            return "Use capability create_presentation (main_title + slides), create_from_outline (outline text), create_from_documents (paths or document_contents), or create_from_source (single source: text or JSON) to generate a PPT. Output is saved to your private output folder and a link is returned."
        except Exception as e:
            logger.exception("ppt-generation run: {}", e)
            return json.dumps({"success": False, "error": str(e)})

    def initialize(self):
        if self.initialized:
            return
        try:
            super().initialize()
            self.initialized = True
        except Exception as e:
            logger.debug("ppt-generation plugin initialize: {}", e)
            self.initialized = True
