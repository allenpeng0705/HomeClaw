#!/usr/bin/env python3
"""
Generate a PowerPoint introduction to HomeClaw.
Requires: pip install python-pptx
Run from repo root:
  python scripts/build_homeclaw_intro_ppt.py           # English → HomeClaw-Intro.pptx
  python scripts/build_homeclaw_intro_ppt.py --lang zh # Chinese → HomeClaw-Intro-zh.pptx

Optional images:
  - Background (all slides): docs/presentations/assets/background.png or background.jpg
  - Logo (title slide): HomeClaw_Banner.jpg in repo root, or docs/presentations/assets/logo.png
  - Architecture (Architecture slide): docs/presentations/assets/system-overview.png
"""
from pathlib import Path
import subprocess
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Install python-pptx first: pip install python-pptx")
    raise

# Layout indices: 0=title, 1=title+content, 6=blank
LAYOUT_TITLE = 0
LAYOUT_TITLE_BODY = 1
LAYOUT_BLANK = 6

# Slide content: en = English, zh = 简体中文
CONTENT = {
    "en": {
        "title_main": "HomeClaw",
        "title_sub": "Your AI assistant. Your hardware. Your control.",
        "s2_title": "What is HomeClaw?",
        "s2_bullets": [
            "An AI assistant that runs on your own hardware",
            "One installation = one autonomous agent",
            "Talks over the channels you already use: WebChat, Telegram, Discord, Email, Companion app",
            "Keeps memory (RAG + agent memory) and extends with plugins and skills",
            "Use cloud models (OpenAI, Gemini, DeepSeek…), local models (llama.cpp, GGUF), or both",
        ],
        "s3_title": "Why HomeClaw?",
        "s3_bullets": [
            "Decentralized & private — data can stay at home",
            "Channel-agnostic — same Core, same memory, any channel",
            "Modular — swap LLM, memory, plugins without changing core logic",
            "Extensible — plugins (any language) and skills (workflows)",
            "Mix mode — smart per-request routing: local vs cloud",
        ],
        "s4_title": "Architecture",
        "s4_bullets": [
            "Layer 1 — Clients: Channels (WebChat, Telegram, Discord, Email) + Companion app",
            "Layer 2 — Core:",
            "  • Memory (RAG + Markdown agent memory)",
            "  • Tools (base for skills)",
            "  • Skills & Plugins (registered, filtered per request)",
            "  • LLM (local and/or cloud)",
        ],
        "s5_title": "Companion app & channels",
        "s5_bullets": [
            "Companion app — Flutter: Mac, Windows, iPhone, Android",
            "  Chat, voice, attachments; Manage Core (edit core.yml, user.yml) from the app",
            "WebChat, CLI, Telegram, Discord, Email — all connect to the same Core",
            "One agent, one memory, one set of tools and plugins",
        ],
        "s6_title": "Memory",
        "s6_bullets": [
            "RAG — vector + relational + optional graph (Cognee default or Chroma)",
            "Agent memory — AGENT_MEMORY.md (long-term), daily memory (short-term)",
            "Per-user context; tools for search and recall",
        ],
        "s7_title": "Mix mode: Smart local/cloud routing",
        "s7_bullets": [
            "Per-request choice: use local or cloud main model",
            "3-layer router (before tools/plugins):",
            "  • Layer 1 — Heuristic (keywords, long-input rules)",
            "  • Layer 2 — Semantic (embedding similarity)",
            "  • Layer 3 — Classifier or Perplexity (small model or main model confidence)",
            "Reports: router decisions and cloud usage via API and usage_report tool",
        ],
        "s8_title": "Plugins & skills",
        "s8_bullets": [
            "Plugins — built-in (Python) and external (any language: Node.js, Go, Java…)",
            "  System plugin example: homeclaw-browser (WebChat UI, browser automation)",
            "Skills — OpenClaw-style workflows (SKILL.md); LLM uses tools and run_skill",
            "Multimodal — images, audio, video (local and cloud models)",
        ],
        "s9_title": "Get started",
        "s9_bullets": [
            "Docs: https://allenpeng0705.github.io/HomeClaw/",
            "Repo: https://github.com/allenpeng0705/HomeClaw",
            "Config: config/core.yml (main_llm, memory, tools, hybrid_router…)",
            "Companion app: clients/HomeClawApp/",
        ],
        "s10_title": "Thank you",
        "s10_sub": "HomeClaw — for the people.\nContributions welcome. Apache 2.0.",
    },
    "zh": {
        "title_main": "HomeClaw",
        "title_sub": "您的 AI 助手。您的硬件。您的掌控。",
        "s2_title": "什么是 HomeClaw？",
        "s2_bullets": [
            "在您自己的硬件上运行的 AI 助手",
            "一次安装 = 一个自主智能体",
            "通过您已在用的渠道对话：WebChat、Telegram、Discord、邮件、伴侣应用",
            "具备记忆（RAG + 智能体记忆），并通过插件与技能扩展能力",
            "可使用云端模型（OpenAI、Gemini、DeepSeek…）、本地模型（llama.cpp、GGUF）或两者兼用",
        ],
        "s3_title": "为什么选 HomeClaw？",
        "s3_bullets": [
            "去中心化与隐私 — 数据可完全留在家中",
            "渠道无关 — 同一 Core、同一记忆、任意渠道",
            "模块化 — 可更换 LLM、记忆、插件，无需改动核心逻辑",
            "可扩展 — 插件（任意语言）与技能（工作流）",
            "Mix 模式 — 按请求智能选择本地或云端",
        ],
        "s4_title": "架构",
        "s4_bullets": [
            "第一层 — 客户端：渠道（WebChat、Telegram、Discord、邮件）+ 伴侣应用",
            "第二层 — Core：",
            "  • 记忆（RAG + Markdown 智能体记忆）",
            "  • 工具（技能基础）",
            "  • 技能与插件（注册在 RAG 中，按请求筛选）",
            "  • LLM（本地和/或云端）",
        ],
        "s5_title": "伴侣应用与渠道",
        "s5_bullets": [
            "伴侣应用 — Flutter：Mac、Windows、iPhone、Android",
            "  聊天、语音、附件；在应用中管理 Core（编辑 core.yml、user.yml）",
            "WebChat、CLI、Telegram、Discord、邮件 — 均连接同一 Core",
            "一个智能体、一套记忆、一套工具与插件",
        ],
        "s6_title": "记忆",
        "s6_bullets": [
            "RAG — 向量 + 关系型 + 可选图（默认 Cognee 或 Chroma）",
            "智能体记忆 — AGENT_MEMORY.md（长期）、每日记忆（短期）",
            "按用户上下文；提供搜索与回忆工具",
        ],
        "s7_title": "Mix 模式：智能本地/云端路由",
        "s7_bullets": [
            "按请求选择：使用本地或云端主模型",
            "三层路由器（在工具/插件之前）：",
            "  • 第一层 — 启发式（关键词、长输入规则）",
            "  • 第二层 — 语义（嵌入相似度）",
            "  • 第三层 — 分类器或困惑度（小模型或主模型置信度）",
            "报告：通过 API 与 usage_report 工具查看路由决策与云端用量",
        ],
        "s8_title": "插件与技能",
        "s8_bullets": [
            "插件 — 内置（Python）与外部（任意语言：Node.js、Go、Java…）",
            "  系统插件示例：homeclaw-browser（WebChat UI、浏览器自动化）",
            "技能 — OpenClaw 风格工作流（SKILL.md）；LLM 使用工具与 run_skill",
            "多模态 — 图像、音频、视频（本地与云端模型）",
        ],
        "s9_title": "快速开始",
        "s9_bullets": [
            "文档：https://allenpeng0705.github.io/HomeClaw/",
            "仓库：https://github.com/allenpeng0705/HomeClaw",
            "配置：config/core.yml（main_llm、memory、tools、hybrid_router…）",
            "伴侣应用：clients/HomeClawApp/",
        ],
        "s10_title": "谢谢",
        "s10_sub": "HomeClaw — 为人而建。\n欢迎贡献。Apache 2.0。",
    },
}


def set_placeholder_text(slide, idx: int, text: str):
    """Set text on a placeholder by index (0=title, 1=body)."""
    slide.placeholders[idx].text = text


def add_bullets(shape, bullets: list[str], level: int = 0):
    """Add bullet paragraphs to a body shape. Level 0 = top-level."""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = level
        p.space_after = Pt(6)


def add_background_image(slide, path: Path, slide_width_in: float, slide_height_in: float):
    """Add a full-slide background image and send it to back. Text stays default (readable)."""
    if not path or not path.exists():
        return
    try:
        pic = slide.shapes.add_picture(
            str(path), Inches(0), Inches(0),
            width=Inches(slide_width_in), height=Inches(slide_height_in)
        )
        sp_tree = pic._element.getparent()
        sp_tree.remove(pic._element)
        sp_tree.insert(0, pic._element)  # first = back in z-order
    except Exception:
        pass


def add_image_if_exists(slide, path: Path, left_in, top_in, width_in=None, height_in=None):
    """Add picture to slide if path exists. Prefer PNG over SVG for compatibility."""
    if not path or not path.exists():
        return None
    try:
        if width_in is not None and height_in is not None:
            return slide.shapes.add_picture(
                str(path), Inches(left_in), Inches(top_in),
                width=Inches(width_in), height=Inches(height_in)
            )
        if height_in is not None:
            return slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in), height=Inches(height_in))
        if width_in is not None:
            return slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in), width=Inches(width_in))
        return slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in))
    except Exception:
        return None


def main():
    lang = "en"
    if "--lang" in sys.argv:
        i = sys.argv.index("--lang")
        if i + 1 < len(sys.argv) and sys.argv[i + 1].lower() in ("zh", "en"):
            lang = sys.argv[i + 1].lower()
    c = CONTENT[lang]

    repo_root = Path(__file__).resolve().parent.parent
    out_dir = repo_root / "docs" / "presentations"
    assets_dir = out_dir / "assets"
    out_dir.mkdir(parents=True, exist_ok=True)
    assets_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / ("HomeClaw-Intro-zh.pptx" if lang == "zh" else "HomeClaw-Intro.pptx")

    # Optional image paths (script works without them)
    logo_candidates = [
        repo_root / "HomeClaw_Banner.jpg",
        assets_dir / "logo.png",
        assets_dir / "logo.jpg",
    ]
    svg_path = repo_root / "docs" / "assets" / "system-overview.svg"
    arch_png = assets_dir / "system-overview.png"
    if not arch_png.exists() and svg_path.exists():
        try:
            import cairosvg
            cairosvg.svg2png(url=str(svg_path), write_to=str(arch_png))
        except Exception:
            pass
        if not arch_png.exists() and sys.platform == "darwin":
            try:
                subprocess.run(
                    ["qlmanage", "-t", "-s", "1200", "-o", str(assets_dir), str(svg_path)],
                    check=True, capture_output=True, timeout=15
                )
                for name in [svg_path.stem + ".png", svg_path.name + ".png", "system-overview.svg.png"]:
                    qlm_out = assets_dir / name
                    if qlm_out.exists():
                        if qlm_out.resolve() != arch_png.resolve():
                            qlm_out.rename(arch_png)
                        break
            except Exception:
                pass
    # Use only PNG for architecture (PowerPoint often doesn't embed SVG from add_picture)
    arch_candidates = [
        arch_png,
        repo_root / "docs" / "assets" / "system-overview.png",
    ]
    logo_path = next((p for p in logo_candidates if p.exists()), None)
    arch_path = next((p for p in arch_candidates if p.exists()), None)
    background_path = next((p for p in [assets_dir / "background.png", assets_dir / "background.jpg"] if p.exists()), None)

    slide_w, slide_h = 10.0, 7.5
    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    slide.shapes.title.text = c["title_main"]
    sub = slide.placeholders[1]
    sub.text = c["title_sub"]
    if logo_path:
        add_image_if_exists(slide, logo_path, 7.2, 0.5, width_in=2.2, height_in=1.0)

    # --- Slide 2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s2_title"])
    add_bullets(slide.placeholders[1], c["s2_bullets"])

    # --- Slide 3 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s3_title"])
    add_bullets(slide.placeholders[1], c["s3_bullets"])

    # --- Slide 4: Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s4_title"])
    add_bullets(slide.placeholders[1], c["s4_bullets"])
    if arch_path:
        add_image_if_exists(slide, arch_path, 5.0, 1.6, width_in=4.6, height_in=5.2)

    # --- Slide 5 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s5_title"])
    add_bullets(slide.placeholders[1], c["s5_bullets"])

    # --- Slide 6 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s6_title"])
    add_bullets(slide.placeholders[1], c["s6_bullets"])

    # --- Slide 7 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s7_title"])
    add_bullets(slide.placeholders[1], c["s7_bullets"])

    # --- Slide 8 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s8_title"])
    add_bullets(slide.placeholders[1], c["s8_bullets"])

    # --- Slide 9 ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    set_placeholder_text(slide, 0, c["s9_title"])
    add_bullets(slide.placeholders[1], c["s9_bullets"])

    # --- Slide 10: Thank you ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    if background_path:
        add_background_image(slide, background_path, slide_w, slide_h)
    slide.shapes.title.text = c["s10_title"]
    slide.placeholders[1].text = c["s10_sub"]

    prs.save(out_path)
    print(f"Saved: {out_path} ({'中文' if lang == 'zh' else 'English'})")
    if not logo_path:
        print("  Tip: Add HomeClaw_Banner.jpg (repo root) or docs/presentations/assets/logo.png for logo on title slide.")
    if not arch_path:
        print("  Tip: Add docs/presentations/assets/system-overview.png (export from docs/assets/system-overview.svg) for architecture slide.")


if __name__ == "__main__":
    main()
