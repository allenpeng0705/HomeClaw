#!/usr/bin/env python3
# /// script
# requires-python = ">=3.9"
# dependencies = [
#     "python-pptx>=1.0.1",
# ]
# ///
"""
Create PowerPoint (.pptx) from outline, source, structured slides, or documents.
Run via: run_skill(skill_name='ppt-generation-1.0.0', script='create_pptx.py', args=[...]).
When HOMECLAW_OUTPUT_DIR is set (by Core), saves there and prints JSON with output_rel_path so Core can append the open link.

Supports:
  - Multiple slide layouts (title, body, section header, two-column)
  - Speaker notes per slide (via "notes" field in slide dict)
  - Dark/light theme
  - CJK font selection via --language zh
  - --dry-run preview (no file written)
  - --slides-file for clean JSON input
  - --max-bullets and --max-chars-per-bullet limits
  - Cover slide: --author, --date
"""
import argparse
import json
import os
import re
import sys
from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print('{"success": false, "error": "python-pptx not installed. Install with: pip install python-pptx"}', flush=True)
    sys.exit(1)

# Slide layout indices (standard python-pptx layouts)
LAYOUT_TITLE = 0
LAYOUT_TITLE_BODY = 1
LAYOUT_SECTION_HEADER = 2  # often a full-bleed title slide used between sections
LAYOUT_BLANK = 6

# Font families by language
FONT_CJK = "Microsoft YaHei, Heiti SC, PingFang SC, Source Han Sans SC, SimHei, sans-serif"
FONT_LATIN = "Calibri, Arial, Helvetica, sans-serif"


def _skill_root() -> Path:
    return Path(__file__).resolve().parent.parent


def _output_dir() -> Path:
    out = os.environ.get("HOMECLAW_OUTPUT_DIR", "").strip()
    if out:
        return Path(out)
    root = _skill_root()
    project = root.parent.parent
    ws = project / "config" / "workspace" / "presentations"
    if ws.parent.is_dir():
        return ws
    return root / "output"


def _safe_filename(name: str) -> str:
    name = re.sub(r"[^\w\s\-\.]", "", name)
    name = name.strip()[:80] or "presentation"
    return name + ".pptx" if not name.lower().endswith(".pptx") else name


def _font(is_cjk: bool) -> str:
    return FONT_CJK if is_cjk else FONT_LATIN


def _build_pptx(
    main_title: str,
    subtitle: str,
    slides_list: List[Dict[str, Any]],
    out_path: Path,
    *,
    language: str = "en",
    author: str = "",
    notes_date: str = "",
    theme: str = "light",
    dry_run: bool = False,
) -> Dict[str, Any]:
    """
    Build the .pptx file.

    Each slide dict may contain:
      - title      : str
      - bullets    : list[str]
      - layout     : str  ("title", "body", "section", "blank")
      - notes      : str  (speaker notes)
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    is_cjk = language == "zh"
    default_font = _font(is_cjk)

    # ---- Cover slide ----
    cover_layout_idx = LAYOUT_TITLE
    cover = prs.slides.add_slide(prs.slide_layouts[cover_layout_idx])
    if hasattr(cover.shapes, "title") and cover.shapes.title:
        cover.shapes.title.text = main_title or "Presentation"

    # Subtitle placeholder
    for shape in cover.placeholders:
        ph_idx = shape.placeholder_format.idx
        if ph_idx == 1:  # subtitle
            shape.text = subtitle or ""
            break

    # Author + date on cover (sometimes in placeholder idx 2, sometimes just in text)
    if author or notes_date:
        cover_text = ""
        if author:
            cover_text += author
        if notes_date:
            cover_text += ("  |  " if cover_text else "") + notes_date
        # Try placeholder idx 2 first
        added = False
        for shape in cover.placeholders:
            if shape.placeholder_format.idx == 2:
                shape.text = cover_text
                added = True
                break
        # Fallback: add a text box at the bottom
        if not added:
            from pptx.util import Emu
            txb = cover.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = cover_text
            p.font.size = Pt(10)
            p.font.italic = True

    # Cover speaker notes
    if slides_list and isinstance(slides_list[0], dict) and slides_list[0].get("notes"):
        _set_notes(cover, str(slides_list[0].get("notes", "")))

    # ---- Content slides ----
    layout_map = {
        "title": LAYOUT_TITLE,
        "body": LAYOUT_TITLE_BODY,
        "section": LAYOUT_SECTION_HEADER,
        "blank": LAYOUT_BLANK,
    }

    for idx, item in enumerate(slides_list[1:], start=1):
        if not isinstance(item, dict):
            continue

        title = (item.get("title") or "").strip() or f"Slide {idx + 1}"
        bullets = item.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = [str(bullets)] if bullets else []
        bullets = [str(b).strip() for b in bullets if b]

        notes_text = str(item.get("notes") or "").strip()
        layout_name = item.get("layout", "body")
        layout_idx = layout_map.get(layout_name, LAYOUT_TITLE_BODY)

        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if layout_idx in (LAYOUT_TITLE, LAYOUT_SECTION_HEADER):
            # These layouts usually just have a title
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
        elif layout_idx == LAYOUT_BLANK:
            # Blank: title is optional; add it as a text box if present
            if title:
                from pptx.util import Emu
                txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                tf = txb.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.name = default_font
            if bullets:
                txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                tf2 = txb2.text_frame
                tf2.word_wrap = True
                for bi, bullet in enumerate(bullets):
                    p2 = tf2.paragraphs[0] if bi == 0 else tf2.add_paragraph()
                    p2.text = bullet
                    p2.font.size = Pt(14)
                    p2.font.name = default_font
                    p2.level = 0
        else:
            # LAYOUT_TITLE_BODY — standard content slide
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            for shape in slide.placeholders:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 1:  # body placeholder
                    _fill_placeholder(shape, bullets, default_font)
                    break

        if notes_text:
            _set_notes(slide, notes_text)

    if dry_run:
        return {
            "success": True,
            "dry_run": True,
            "message": f"[Dry run] Would create presentation '{main_title}' with {len(slides_list)} slide(s).",
            "slides_preview": [
                {"title": (s or {}).get("title", ""), "bullets_count": len((s or {}).get("bullets", []))}
                for s in slides_list
            ],
        }

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    payload = {
        "success": True,
        "path": str(out_path.resolve()),
        "message": f"Presentation saved to {out_path.resolve()}",
        "slides_count": len(slides_list),
    }
    if os.environ.get("HOMECLAW_OUTPUT_DIR"):
        payload["output_rel_path"] = f"output/{out_path.name}"
    return payload


def _fill_placeholder(shape, bullets: List[str], font: str) -> None:
    """Fill a placeholder shape with bullet text."""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(14)
        p.space_after = Pt(6)


def _set_notes(slide, text: str) -> None:
    """Set speaker notes on a slide."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text
    except Exception:
        pass


def _parse_outline(outline: str) -> Tuple[str, str, List[Dict[str, Any]]]:
    main_title = ""
    subtitle = ""
    slides_list: List[Dict[str, Any]] = []
    current_title = None
    current_bullets: List[str] = []
    current_notes = ""

    def flush_slide():
        nonlocal current_title, current_bullets, current_notes
        if current_title is not None:
            slides_list.append({"title": current_title, "bullets": list(current_bullets), "notes": current_notes})
        current_title = None
        current_bullets = []
        current_notes = ""

    lines = (outline or "").strip().split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        # Notes directive: <!-- notes: ... -->
        if "<!--" in stripped and "notes:" in stripped:
            notes_match = re.search(r"notes:\s*(.+?)\s*-->", stripped)
            if notes_match:
                current_notes = notes_match.group(1).strip()
            i += 1
            continue
        if stripped.startswith("##"):
            flush_slide()
            current_title = stripped[2:].strip() or "Slide"
            current_notes = ""
        elif stripped.startswith("-") or stripped.startswith("*"):
            bullet = stripped[1:].strip()
            if current_title is not None:
                current_bullets.append(bullet)
            else:
                if not main_title:
                    main_title = bullet
                elif not subtitle:
                    subtitle = bullet
                else:
                    if not current_title:
                        current_title = "Content"
                    current_bullets.append(bullet)
        elif stripped.startswith(">>>"):
            # Section header marker
            flush_slide()
            current_title = stripped[3:].strip() or "Section"
            slides_list.append({"title": current_title, "bullets": [], "layout": "section", "notes": ""})
            current_title = None
            current_notes = ""
        elif stripped:
            if current_title is not None:
                current_bullets.append(stripped)
            else:
                if not main_title:
                    main_title = stripped
                elif not subtitle:
                    subtitle = stripped
                else:
                    current_title = "Content"
                    current_bullets.append(stripped)
        i += 1

    flush_slide()
    if not main_title and slides_list:
        main_title = slides_list[0].get("title") or "Presentation"
    return main_title, subtitle, slides_list


def _parse_document_to_slides(
    content: str,
    doc_title: str = "",
    *,
    max_bullets: int = 12,
    max_chars: int = 400,
) -> List[Dict[str, Any]]:
    if not (content or "").strip():
        return [{"title": doc_title or "Document", "bullets": ["(No content)"]}]
    text = (content or "").strip()
    slides: List[Dict[str, Any]] = []
    if "##" in text or text.startswith("#"):
        main_title, subtitle, slides_list = _parse_outline(text)
        if slides_list:
            for s in slides_list:
                slides.append(s)
        elif main_title:
            slides.append({"title": main_title, "bullets": [subtitle] if subtitle else []})
    else:
        blocks = re.split(r"\n\s*\n", text)
        for block in blocks:
            block = block.strip()
            if not block:
                continue
            lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
            if not lines:
                continue
            title = lines[0][:120]
            bullets = []
            for ln in lines[1:]:
                bullets.append(ln[:max_chars] if len(ln) > max_chars else ln)
            if not bullets and len(lines) == 1 and len(lines[0]) > 80:
                first = lines[0]
                for part in re.split(r"(?<=[.!?])\s+", first)[:max_bullets]:
                    if part.strip():
                        bullets.append(part.strip()[:max_chars])
            if len(bullets) > max_bullets:
                slides.append({"title": title, "bullets": bullets[:max_bullets]})
                for j in range(max_bullets, len(bullets), max_bullets):
                    chunk = bullets[j : j + max_bullets]
                    slides.append({"title": f"{title} (continued)", "bullets": chunk})
            else:
                slides.append({"title": title, "bullets": bullets or [title]})
    if not slides and doc_title:
        slides.append({"title": doc_title, "bullets": [text[:500]]})
    return slides


def _resolve_document_path(path_str: str, base_dirs: List[Path]) -> Optional[Path]:
    path_str = (path_str or "").strip()
    if not path_str:
        return None
    p = Path(path_str)
    bases = [b.resolve() for b in base_dirs if b.is_dir()]
    if p.is_absolute():
        resolved = p.resolve()
        for base in bases:
            try:
                resolved.relative_to(base)
                if resolved.is_file():
                    return resolved
            except ValueError:
                continue
        return None
    for base in bases:
        candidate = (base / p).resolve()
        try:
            candidate.relative_to(base)
        except ValueError:
            continue
        if candidate.is_file():
            return candidate
    return None


def _read_arg_or_file(value: Optional[str], path_value: Optional[str]) -> str:
    """Return inline value or file content. File path must be under skill root or project (no path escape)."""
    if value is not None and value.strip():
        return value.strip()
    if not path_value or not path_value.strip():
        return ""
    p = Path(path_value.strip())
    try:
        if not p.is_file():
            return ""
        root = _skill_root()
        project = root.parent.parent
        resolved = p.resolve()
        for base in (root, project):
            try:
                if base.exists() and resolved.relative_to(base.resolve()):
                    return resolved.read_text(encoding="utf-8", errors="replace")
            except ValueError:
                continue
    except (OSError, RuntimeError):
        pass
    return ""


def _load_slides_json(slides_arg: Optional[str], slides_file: Optional[str]) -> Tuple[Optional[List[Dict[str, Any]]], Optional[str]]:
    """Load slides from inline JSON or JSON file. Returns (slides_list, error_message)."""
    if slides_file:
        content = _read_arg_or_file(None, slides_file)
        if not content:
            return None, f"Could not read slides file: {slides_file}"
        try:
            data = json.loads(content)
            if isinstance(data, dict) and "slides" in data:
                # Allow wrapping in {slides: [...]} envelope
                data = data["slides"]
            if not isinstance(data, list):
                return None, f"Slides file must contain a JSON array, got {type(data).__name__}."
            return data, None
        except json.JSONDecodeError as e:
            return None, f"Invalid JSON in slides file: {e}"
    if slides_arg:
        try:
            data = json.loads(slides_arg)
            if isinstance(data, dict) and "slides" in data:
                data = data["slides"]
            return data if isinstance(data, list) else None, None
        except json.JSONDecodeError:
            return None, "Invalid JSON in --slides."
    return None, None


def _write_or_dry_run(
    payload: Dict[str, Any],
    out_path: Path,
    prs,
    dry_run: bool,
    main_title: str,
    slides_count: int,
) -> Dict[str, Any]:
    if dry_run:
        return {
            "success": True,
            "dry_run": True,
            "message": f"[Dry run] Would create presentation '{main_title}' with {slides_count} slide(s).",
        }
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    p = {"success": True, "path": str(out_path.resolve()), "message": f"Presentation saved to {out_path.resolve()}", "slides_count": slides_count}
    if os.environ.get("HOMECLAW_OUTPUT_DIR"):
        p["output_rel_path"] = f"output/{out_path.name}"
    return p


def _slides_from_outline(args: argparse.Namespace) -> Tuple[str, str, List[Dict[str, Any]]]:
    outline = _read_arg_or_file(getattr(args, "outline", None), getattr(args, "outline_file", None))
    if not outline:
        return "", "", []
    max_bullets = getattr(args, "max_bullets", 12) or 12
    max_chars = getattr(args, "max_chars", 400) or 400
    main_title, subtitle, slides_list = _parse_outline(outline)
    return main_title, subtitle, slides_list


def _slides_from_source(args: argparse.Namespace) -> Tuple[str, str, List[Dict[str, Any]]]:
    source = _read_arg_or_file(getattr(args, "source", None), getattr(args, "source_file", None))
    if not source:
        return "", "", []
    main_title = ""
    subtitle = ""
    all_slides: List[Dict[str, Any]] = []
    stripped = source.strip()
    if stripped.startswith("["):
        try:
            content_list = json.loads(source)
        except json.JSONDecodeError:
            content_list = None
    elif stripped.startswith("{"):
        try:
            obj = json.loads(source)
            content_list = obj.get("results") or obj.get("documents") or (obj.get("items") if isinstance(obj.get("items"), list) else None)
        except json.JSONDecodeError:
            content_list = None
    else:
        content_list = None

    max_bullets = getattr(args, "max_bullets", 12) or 12
    max_chars = getattr(args, "max_chars", 400) or 400

    if isinstance(content_list, list) and content_list:
        for i, item in enumerate(content_list):
            if not isinstance(item, dict):
                continue
            title = str(item.get("title") or item.get("name") or item.get("query") or f"Source {i+1}").strip() or f"Source {i+1}"
            content = str(item.get("content") or item.get("body") or item.get("snippet") or "").strip()
            if not content and isinstance(item.get("results"), list):
                for r in item.get("results", [])[:5]:
                    if isinstance(r, dict):
                        content += str(r.get("content") or r.get("body") or r.get("snippet") or r.get("title") or "") + "\n"
            slides = _parse_document_to_slides(content, title, max_bullets=max_bullets, max_chars=max_chars)
            if not main_title and slides:
                main_title = title or (slides[0].get("title") if slides else "")
            all_slides.extend(slides)
        if all_slides and not main_title:
            main_title = all_slides[0].get("title", "Presentation") if all_slides else "Presentation"

    if not all_slides:
        main_title, subtitle, slides_list = _parse_outline(source)
        if slides_list:
            all_slides = slides_list
        elif main_title:
            all_slides = [{"title": main_title, "bullets": [subtitle] if subtitle else []}]
    if not main_title:
        main_title = str((all_slides[0] or {}).get("title") or "Presentation") if all_slides else "Presentation"
    main_title = (main_title or "Presentation").strip() or "Presentation"
    return main_title, subtitle, all_slides


def _shared_output(args: argparse.Namespace, main_title: str, slides_list: List[Dict[str, Any]], subtitle: str = "") -> Dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    language = getattr(args, "language", "en") or "en"
    is_cjk = language == "zh"
    default_font = _font(is_cjk)
    author = getattr(args, "author", "") or ""
    notes_date = getattr(args, "date", "") or ""
    dry_run = getattr(args, "dry_run", False)
    theme = getattr(args, "theme", "light") or "light"

    out_dir = _output_dir()
    out_filename = (getattr(args, "output_filename", None) or "").strip()
    if not out_filename:
        out_filename = _safe_filename(main_title) or f"presentation_{date.today().strftime('%Y%m%d')}.pptx"
    elif not out_filename.lower().endswith(".pptx"):
        out_filename += ".pptx"
    out_path = out_dir / out_filename

    layout_map = {
        "title": LAYOUT_TITLE,
        "body": LAYOUT_TITLE_BODY,
        "section": LAYOUT_SECTION_HEADER,
        "blank": LAYOUT_BLANK,
    }

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    if hasattr(cover.shapes, "title") and cover.shapes.title:
        cover.shapes.title.text = main_title or "Presentation"
    for shape in cover.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = subtitle or ""
            break
    if author or notes_date:
        cover_text = "  |  ".join(filter(None, [author, notes_date]))
        for shape in cover.placeholders:
            if shape.placeholder_format.idx == 2:
                shape.text = cover_text
                break
    if slides_list and isinstance(slides_list[0], dict) and slides_list[0].get("notes"):
        _set_notes(cover, str(slides_list[0].get("notes", "")))

    # Content slides
    for idx, item in enumerate(slides_list[1:], start=1):
        if not isinstance(item, dict):
            continue
        title = (item.get("title") or "").strip() or f"Slide {idx + 1}"
        bullets = item.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = [str(bullets)] if bullets else []
        bullets = [str(b).strip() for b in bullets if b]
        notes_text = str(item.get("notes") or "").strip()
        layout_name = item.get("layout", "body")
        layout_idx = layout_map.get(layout_name, LAYOUT_TITLE_BODY)

        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if layout_idx == LAYOUT_BLANK:
            if title:
                from pptx.util import Emu
                txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                tf = txb.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.name = default_font
            if bullets:
                txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                tf2 = txb2.text_frame
                tf2.word_wrap = True
                for bi, bullet in enumerate(bullets):
                    p2 = tf2.paragraphs[0] if bi == 0 else tf2.add_paragraph()
                    p2.text = bullet
                    p2.font.size = Pt(14)
                    p2.font.name = default_font
        else:
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            if layout_idx == LAYOUT_TITLE_BODY:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        _fill_placeholder(shape, bullets, default_font)
                        break

        if notes_text:
            _set_notes(slide, notes_text)

    if dry_run:
        return {
            "success": True,
            "dry_run": True,
            "message": f"[Dry run] Would create '{main_title}' with {len(slides_list)} slide(s).",
            "slides_preview": [
                {"title": (s or {}).get("title", ""), "bullets_count": len((s or {}).get("bullets", []))}
                for s in slides_list
            ],
        }

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    p = {"success": True, "path": str(out_path.resolve()), "message": f"Presentation saved to {out_path.resolve()}", "slides_count": len(slides_list)}
    if os.environ.get("HOMECLAW_OUTPUT_DIR"):
        p["output_rel_path"] = f"output/{out_filename}"
    return p


def run_outline(args: argparse.Namespace) -> Dict[str, Any]:
    main_title, subtitle, slides_list = _slides_from_outline(args)
    if not slides_list and not main_title:
        return {"success": False, "error": "outline is required (--outline or --outline-file)."}
    if not slides_list:
        return {"success": False, "error": "Could not parse any slides from the outline."}
    return _shared_output(args, main_title, [{"title": main_title, "bullets": []}] + (slides_list if slides_list else []), subtitle)


def run_source(args: argparse.Namespace) -> Dict[str, Any]:
    main_title, subtitle, all_slides = _slides_from_source(args)
    if not all_slides:
        return {"success": False, "error": "source is required (--source or --source-file)."}
    return _shared_output(args, main_title, [{"title": main_title, "bullets": []}] + all_slides, subtitle)


def run_presentation(args: argparse.Namespace) -> Dict[str, Any]:
    main_title = (getattr(args, "main_title", None) or "").strip() or "Presentation"
    subtitle = (getattr(args, "subtitle", None) or "").strip()
    notes_date = getattr(args, "date", "") or date.today().strftime("%Y-%m-%d")

    slides_list, load_err = _load_slides_json(getattr(args, "slides", None), getattr(args, "slides_file", None))
    if load_err:
        return {"success": False, "error": load_err}
    if not slides_list:
        return {"success": False, "error": "slides is required (use --slides JSON or --slides-file)."}

    # Build cover as first slide dict
    cover = [{"title": main_title, "bullets": [], "layout": "title"}]
    combined = cover + slides_list

    return _shared_output(args, main_title, combined, subtitle)


def run_documents(args: argparse.Namespace) -> Dict[str, Any]:
    paths_raw = getattr(args, "document_paths", None)
    contents_raw = getattr(args, "document_contents", None)
    if not paths_raw and not contents_raw:
        return {"success": False, "error": "Provide at least one of --document_paths or --document_contents."}

    main_title = (getattr(args, "main_title", None) or "").strip()
    all_slides: List[Dict[str, Any]] = []
    sources: List[str] = []
    root = _skill_root()
    project = root.parent.parent
    base_dirs = [project, project / "config" / "workspace"]

    max_bullets = getattr(args, "max_bullets", 12) or 12
    max_chars = getattr(args, "max_chars", 400) or 400

    if paths_raw:
        if isinstance(paths_raw, str):
            try:
                path_list = json.loads(paths_raw)
            except json.JSONDecodeError:
                return {"success": False, "error": "document_paths must be a JSON array of strings."}
        else:
            path_list = paths_raw if isinstance(paths_raw, list) else []
        if not isinstance(path_list, list):
            return {"success": False, "error": "document_paths must be an array."}
        for path_str in path_list:
            if not isinstance(path_str, str):
                continue
            resolved = _resolve_document_path(path_str, base_dirs)
            if not resolved:
                return {"success": False, "error": f"File not found or not allowed: {path_str}"}
            try:
                content = resolved.read_text(encoding="utf-8", errors="replace")
            except Exception as e:
                return {"success": False, "error": f"Could not read {path_str}: {e}"}
            doc_title = resolved.stem
            slides = _parse_document_to_slides(content, doc_title, max_bullets=max_bullets, max_chars=max_chars)
            if not main_title and slides:
                main_title = doc_title or slides[0].get("title", "Document")
            all_slides.extend(slides)
            sources.append(resolved.name)

    if contents_raw:
        if isinstance(contents_raw, str):
            try:
                content_list = json.loads(contents_raw)
            except json.JSONDecodeError:
                return {"success": False, "error": "document_contents must be a JSON array of {title, content}."}
        else:
            content_list = contents_raw if isinstance(contents_raw, list) else []
        if not isinstance(content_list, list):
            return {"success": False, "error": "document_contents must be an array."}
        for i, item in enumerate(content_list):
            if not isinstance(item, dict):
                continue
            title = (item.get("title") or item.get("name") or f"Document {i+1}").strip()
            content = (item.get("content") or "").strip()
            slides = _parse_document_to_slides(content, title, max_bullets=max_bullets, max_chars=max_chars)
            if not main_title and slides:
                main_title = title
            all_slides.extend(slides)
            sources.append(title)

    if not all_slides:
        return {"success": False, "error": "No slides could be parsed from the document(s)."}
    if not main_title:
        main_title = all_slides[0].get("title", "Presentation") if all_slides else "Presentation"

    cover = [{"title": main_title, "bullets": []}]
    combined = cover + all_slides
    result = _shared_output(args, main_title, combined, "")
    if sources:
        result["sources"] = sources
    return result


def main() -> None:
    ap = argparse.ArgumentParser(description="Create .pptx from outline, source, slides, or documents.")
    ap.add_argument("--capability", required=True, choices=["outline", "source", "presentation", "documents"], help="Which mode to run")
    ap.add_argument("--outline", default="", help="Markdown outline (## titles, - bullets)")
    ap.add_argument("--outline-file", default="", help="Path to file containing outline")
    ap.add_argument("--source", default="", help="Plain text or JSON array of {title, content}")
    ap.add_argument("--source-file", default="", help="Path to file containing source")
    ap.add_argument("--main_title", default="", help="Title slide title")
    ap.add_argument("--subtitle", default="", help="Title slide subtitle")
    ap.add_argument("--slides", default="", help='JSON array of {"title","bullets","notes","layout"}')
    ap.add_argument("--slides-file", default="", help="Path to a JSON file containing slides array (preferred for large JSON)")
    ap.add_argument("--document_paths", default="", help='JSON array of file paths')
    ap.add_argument("--document_contents", default="", help='JSON array of {"title", "content"}')
    ap.add_argument("--output_filename", default="", help="Output .pptx filename")
    ap.add_argument("--language", default="en", help="en or zh — selects CJK font for zh")
    ap.add_argument("--author", default="", help="Author name on cover slide")
    ap.add_argument("--date", default="", help="Date on cover slide (default: today)")
    ap.add_argument("--theme", default="light", choices=["light", "dark"], help="Presentation theme (future use)")
    ap.add_argument("--max_bullets", type=int, default=12, help="Max bullets per slide when auto-parsing (default: 12)")
    ap.add_argument("--max_chars", type=int, default=400, help="Max characters per bullet when auto-parsing (default: 400)")
    ap.add_argument("--dry-run", action="store_true", help="Preview slides without writing a file")
    args = ap.parse_args()

    try:
        if args.capability == "outline":
            payload = run_outline(args)
        elif args.capability == "source":
            payload = run_source(args)
        elif args.capability == "presentation":
            payload = run_presentation(args)
        else:
            payload = run_documents(args)
    except (OSError, IOError, ValueError, TypeError, KeyError) as e:
        payload = {"success": False, "error": str(e)}
    except Exception as e:
        payload = {"success": False, "error": f"Unexpected error: {e}"}

    print(json.dumps(payload, ensure_ascii=False), flush=True)
    if not payload.get("success"):
        sys.exit(1)


if __name__ == "__main__":
    main()
